from __future__ import annotations

import csv
import hashlib
import json
import re
import uuid
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config.settings import Settings, get_settings


@dataclass(slots=True)
class LoadedDocument:
    document_id: str
    source_name: str
    source_path: str
    text: str
    metadata: dict[str, Any]


@dataclass(slots=True)
class DocumentChunk:
    chunk_id: str
    document_id: str
    source_name: str
    text: str
    metadata: dict[str, Any]


class DocumentIngestor:
    """Load diverse document types and split them into indexed chunks."""

    def __init__(self, settings: Settings | None = None) -> None:
        self.settings = settings or get_settings()
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.settings.chunk_size,
            chunk_overlap=self.settings.chunk_overlap,
            separators=["\n\n", "\n", " ", ""],
        )

    def load_files(self, paths: list[Path | str]) -> list[LoadedDocument]:
        loaded_documents: list[LoadedDocument] = []
        for path in paths:
            loaded_documents.extend(self.load_file(Path(path)))
        return loaded_documents

    def load_file(self, path: Path) -> list[LoadedDocument]:
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {path}")

        suffix = path.suffix.lower()
        document_id = hashlib.sha256(path.read_bytes()).hexdigest()
        loader = self._get_loader(suffix)
        return loader(path, document_id)

    def chunk_documents(self, documents: list[LoadedDocument]) -> list[DocumentChunk]:
        chunks: list[DocumentChunk] = []
        for document in documents:
            segments = self.splitter.split_text(self._normalize_text(document.text))
            for index, text in enumerate(segments):
                if self.settings.enable_deduplication:
                    chunk_name = f"{document.document_id}:{index}"
                    chunk_id = str(uuid.uuid5(uuid.NAMESPACE_OID, chunk_name))
                else:
                    chunk_id = uuid.uuid4().hex
                chunks.append(
                    DocumentChunk(
                        chunk_id=chunk_id,
                        document_id=document.document_id,
                        source_name=document.source_name,
                        text=text,
                        metadata={**document.metadata, "chunk_index": index},
                    )
                )
        return chunks

    def supported_extensions(self) -> list[str]:
        return list(self.settings.allowed_extensions)

    def _get_loader(self, suffix: str):
        loaders = {
            ".txt": self._load_text,
            ".md": self._load_text,
            ".rst": self._load_text,
            ".rtf": self._load_text,
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".pptx": self._load_pptx,
            ".html": self._load_html,
            ".htm": self._load_html,
            ".csv": self._load_csv,
            ".json": self._load_json,
            ".xml": self._load_xml,
            ".png": self._load_image,
            ".jpg": self._load_image,
            ".jpeg": self._load_image,
            ".tif": self._load_image,
            ".tiff": self._load_image,
        }

        if suffix not in loaders:
            raise ValueError(f"Unsupported file format: {suffix}")
        return loaders[suffix]

    def _load_text(self, path: Path, document_id: str) -> list[LoadedDocument]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "text/plain"})]

    def _load_pdf(self, path: Path, document_id: str) -> list[LoadedDocument]:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        loaded_documents: list[LoadedDocument] = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if not text.strip():
                text = self._extract_pdf_page_text_fallback(path, page_number)
            loaded_documents.append(
                self._build_loaded_document(
                    path,
                    document_id,
                    text,
                    {"mime_type": "application/pdf", "page": page_number},
                )
            )
        return loaded_documents

    def _extract_pdf_page_text_fallback(self, path: Path, page_number: int) -> str:
        try:
            import fitz
        except ImportError:
            return ""

        try:
            with fitz.open(str(path)) as pdf:
                if 1 <= page_number <= pdf.page_count:
                    return pdf.load_page(page_number - 1).get_text("text") or ""
        except Exception:
            return ""

        return ""

    def _load_docx(self, path: Path, document_id: str) -> list[LoadedDocument]:
        from docx import Document

        document = Document(str(path))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"})]

    def _load_pptx(self, path: Path, document_id: str) -> list[LoadedDocument]:
        from pptx import Presentation

        presentation = Presentation(str(path))
        slide_text: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
            if slide_parts:
                slide_text.append(f"Slide {slide_index}: " + " ".join(slide_parts))
        text = "\n".join(slide_text)
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})]

    def _load_html(self, path: Path, document_id: str) -> list[LoadedDocument]:
        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        text = soup.get_text(" ", strip=True)
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "text/html"})]

    def _load_csv(self, path: Path, document_id: str) -> list[LoadedDocument]:
        rows: list[str] = []
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
            reader = csv.reader(handle)
            for row in reader:
                rows.append(" | ".join(cell.strip() for cell in row))
        return [self._build_loaded_document(path, document_id, "\n".join(rows), {"mime_type": "text/csv"})]

    def _load_json(self, path: Path, document_id: str) -> list[LoadedDocument]:
        payload = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        text = json.dumps(payload, indent=2, ensure_ascii=False)
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "application/json"})]

    def _load_xml(self, path: Path, document_id: str) -> list[LoadedDocument]:
        root = ET.fromstring(path.read_text(encoding="utf-8", errors="ignore"))
        text = ET.tostring(root, encoding="unicode", method="xml")
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "application/xml"})]

    def _load_image(self, path: Path, document_id: str) -> list[LoadedDocument]:
        if not self.settings.enable_ocr:
            raise ValueError(f"OCR is disabled for image input: {path.name}")

        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:  # pragma: no cover - optional OCR dependency
            raise ValueError("Image OCR requires pytesseract and Pillow") from exc

        text = pytesseract.image_to_string(Image.open(path))
        return [self._build_loaded_document(path, document_id, text, {"mime_type": "image/*", "ocr": True})]

    def _build_loaded_document(
        self,
        path: Path,
        document_id: str,
        text: str,
        extra_metadata: dict[str, Any],
    ) -> LoadedDocument:
        metadata = {
            "source_name": path.name,
            "source_path": str(path),
            "document_id": document_id,
            **extra_metadata,
        }
        return LoadedDocument(
            document_id=document_id,
            source_name=path.name,
            source_path=str(path),
            text=text,
            metadata=metadata,
        )

    @staticmethod
    def _normalize_text(text: str) -> str:
        text = re.sub(r"\r\n?", "\n", text)
        text = re.sub(r"[ \t]+", " ", text)
        return text.strip()