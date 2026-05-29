from __future__ import annotations

from csv import writer
from json import dumps
from pathlib import Path

from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
TARGET = ROOT / "data" / "test_corpus"


def make_repeated_text(title: str, size_kb: int) -> str:
    paragraph = (
        f"{title}: This corpus is designed for retrieval testing, citation traceability, and production showcase demos. "
        "It contains repeated operational language about indexing, retrieval quality, observability, and safety controls. "
    )
    chunks = []
    while sum(len(chunk) for chunk in chunks) < size_kb * 1024:
        chunks.append(paragraph)
    return "\n\n".join(chunks)


def write_txt(path: Path, title: str, size_kb: int) -> None:
    path.write_text(make_repeated_text(title, size_kb), encoding="utf-8")


def write_md(path: Path, title: str, size_kb: int) -> None:
    body = "# " + title + "\n\n" + make_repeated_text(title, size_kb)
    path.write_text(body, encoding="utf-8")


def write_html(path: Path, title: str, size_kb: int) -> None:
    body = make_repeated_text(title, size_kb)
    path.write_text(f"<html><body><h1>{title}</h1><p>{body}</p></body></html>", encoding="utf-8")


def write_json(path: Path, title: str, size_kb: int) -> None:
    payload = {"title": title, "content": make_repeated_text(title, size_kb)}
    path.write_text(dumps(payload, indent=2), encoding="utf-8")


def write_csv(path: Path, title: str, size_kb: int) -> None:
    rows = [["section", "text"]]
    body = make_repeated_text(title, size_kb)
    for index, segment in enumerate(body.split("\n\n"), start=1):
        rows.append([f"{title}-{index}", segment])

    with path.open("w", encoding="utf-8", newline="") as handle:
        csv_writer = writer(handle)
        csv_writer.writerows(rows)


def write_xml(path: Path, title: str, size_kb: int) -> None:
    body = make_repeated_text(title, size_kb)
    path.write_text(f"<document><title>{title}</title><body>{body}</body></document>", encoding="utf-8")


def write_docx(path: Path, title: str, size_kb: int) -> None:
    document = Document()
    document.add_heading(title, 0)
    for paragraph in make_repeated_text(title, size_kb).split("\n\n"):
        document.add_paragraph(paragraph)
    document.save(path)


def write_pptx(path: Path, title: str, size_kb: int) -> None:
    presentation = Presentation()
    body = make_repeated_text(title, size_kb).split("\n\n")
    for index, paragraph in enumerate(body[:10], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{title} - Slide {index}"
        textbox = slide.shapes.add_textbox(0, 1000000, 8000000, 4000000)
        textbox.text_frame.text = paragraph
    presentation.save(path)


def write_pdf(path: Path, title: str, size_kb: int) -> None:
    body = make_repeated_text(title, size_kb)
    pdf = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    pdf.setTitle(title)
    y = height - 50
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(40, y, title)
    pdf.setFont("Helvetica", 10)
    for line in body.split("\n"):
        if y < 60:
            pdf.showPage()
            y = height - 50
            pdf.setFont("Helvetica", 10)
        pdf.drawString(40, y, line[:120])
        y -= 12
    pdf.save()


def main() -> None:
    TARGET.mkdir(parents=True, exist_ok=True)
    specs = [
        ("company_overview.txt", write_txt, "Company Overview", 1500),
        ("product_notes.md", write_md, "Product Notes", 1500),
        ("ops_handbook.html", write_html, "Operations Handbook", 1500),
        ("metrics.json", write_json, "Metrics Snapshot", 1500),
        ("support.csv", write_csv, "Support Matrix", 1500),
        ("taxonomy.xml", write_xml, "Taxonomy", 1500),
        ("research_brief.docx", write_docx, "Research Brief", 1500),
        ("executive_deck.pptx", write_pptx, "Executive Deck", 1500),
        ("architecture.pdf", write_pdf, "Architecture Notes", 1500),
        ("release_plan.txt", write_txt, "Release Plan", 1500),
    ]

    for file_name, writer_fn, title, size_kb in specs:
        writer_fn(TARGET / file_name, title, size_kb)

    print(f"Generated {len(specs)} documents in {TARGET}")


if __name__ == "__main__":
    main()
